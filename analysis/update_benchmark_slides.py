from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DECK = ROOT / "slides" / "prediction_results_deck.pptx"
VAL_FIG = ROOT / "plots" / "validation_benchmark_mirrored.png"
LEARN_FIG = ROOT / "plots" / "learning_r2_benchmark.png"

BG = RGBColor(247, 245, 242)
TEXT = RGBColor(31, 41, 51)
MUTED = RGBColor(90, 102, 114)
NAVY = RGBColor(22, 50, 79)
TEAL = RGBColor(42, 157, 143)
PANEL = RGBColor(255, 255, 255)


def clear_slide(slide) -> None:
    for shape in list(slide.shapes):
        sp = shape.element
        sp.getparent().remove(sp)


def add_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(11.9), Inches(0.65))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = NAVY

    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.55),
        Inches(0.98),
        Inches(12.05),
        Inches(0.06),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.05), Inches(11.8), Inches(0.3))
        p = sub_box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = "Aptos"
        r.font.size = Pt(10.5)
        r.font.color.rgb = MUTED


def add_panel(slide, x: float, y: float, w: float, h: float) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = RGBColor(222, 226, 230)


def add_bullets(slide, lines: list[str], x: float, y: float, w: float, h: float) -> None:
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(253, 250, 245)
    panel.line.color.rgb = RGBColor(235, 221, 208)

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(0.12),
        Inches(h),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()

    box = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.15), Inches(w - 0.32), Inches(h - 0.2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Aptos"
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT
        p.bullet = True
        p.space_after = Pt(5)


def add_takeaway(slide, text: str) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.65),
        Inches(6.72),
        Inches(12.0),
        Inches(0.42),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.85), Inches(6.79), Inches(11.7), Inches(0.24))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER


def rebuild_slide(slide, title: str, subtitle: str, fig_path: Path, bullets: list[str], takeaway: str) -> None:
    clear_slide(slide)
    add_background(slide)
    add_title(slide, title, subtitle)
    add_panel(slide, 0.55, 1.22, 8.4, 5.35)
    slide.shapes.add_picture(str(fig_path), Inches(0.7), Inches(1.35), width=Inches(8.1))
    add_bullets(slide, bullets, 9.05, 1.38, 3.5, 4.7)
    add_takeaway(slide, takeaway)


def main() -> None:
    prs = Presentation(DECK)

    rebuild_slide(
        prs.slides[1],
        "Benchmark Before Augmentation (Validation Wave)",
        "Baseline-only comparison across elicitation modes; elastic net and the noise ceiling are kept in the slide text.",
        VAL_FIG,
        [
            "GPT-4.1 baseline_reasoning is the strongest GPT-4.1 baseline on validation: RMSE 5.70, correlation 0.627, directional accuracy 0.70, R² 0.362.",
            "GPT-4.1-nano is weaker across all four metrics on validation; nano baseline_reasoning drops to R² -0.994.",
            "External benchmark: elastic net reaches RMSE 4.52, correlation 0.773, directional accuracy 0.75, R² 0.599; the validation noise ceiling is R² 0.641.",
        ],
        "Validation is the cleanest place to show that elastic net still outperforms the non-augmented LLM baselines.",
    )

    rebuild_slide(
        prs.slides[2],
        "Benchmark Before Augmentation (Learning Wave)",
        "Baseline-only comparison on the learning wave, using the same four metrics and the same layout as validation.",
        LEARN_FIG,
        [
            "GPT-4.1 is fairly flat across elicitation modes on learning; baseline_reasoning is best on RMSE and R².",
            "GPT-4.1-nano behaves differently: joint_reasoning is the strongest nano baseline on RMSE, R², and directional accuracy.",
            "Unlike validation, the learning-wave ordering is not dominated by a single clear baseline family across both models.",
        ],
        "The learning wave changes the baseline ordering, especially for GPT-4.1-nano.",
    )

    prs.save(DECK)
    print(DECK)


if __name__ == "__main__":
    main()
