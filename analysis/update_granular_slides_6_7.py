from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DECK = ROOT / "slides" / "prediction_results_deck.pptx"
from plot_paths import GRANULAR_PLOTS

FIG_PERF = GRANULAR_PLOTS / "granular_performance_delta_r2_heatmap.png"
FIG_CHAT = GRANULAR_PLOTS / "granular_chat_shift_heatmap.png"
FIG_ENET = GRANULAR_PLOTS / "llm_shift_vs_enet_importance.png"

BG = RGBColor(247, 245, 242)
TEXT = RGBColor(31, 41, 51)
MUTED = RGBColor(90, 102, 114)
NAVY = RGBColor(22, 50, 79)
TEAL = RGBColor(42, 157, 143)
PANEL = RGBColor(255, 255, 255)


def clear_slide(slide) -> None:
    for shape in list(slide.shapes):
        sp = shape.element
        sp.getparent().remove(sp)


def add_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(11.9), Inches(0.65))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = NAVY

    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.55),
        Inches(0.98),
        Inches(12.05),
        Inches(0.06),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.05), Inches(11.8), Inches(0.3))
        p = sub_box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = "Aptos"
        r.font.size = Pt(10.5)
        r.font.color.rgb = MUTED


def add_panel(slide, x: float, y: float, w: float, h: float) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = RGBColor(222, 226, 230)


def add_bullets(slide, lines: list[str], x: float, y: float, w: float, h: float) -> None:
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(253, 250, 245)
    panel.line.color.rgb = RGBColor(235, 221, 208)

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(0.12),
        Inches(h),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()

    box = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.15), Inches(w - 0.32), Inches(h - 0.2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Aptos"
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT
        p.bullet = True
        p.space_after = Pt(4)


def add_takeaway(slide, text: str) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.65),
        Inches(6.72),
        Inches(12.0),
        Inches(0.42),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.85), Inches(6.79), Inches(11.7), Inches(0.24))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER


def rebuild_slide(slide, title: str, subtitle: str, fig_path: Path, bullets: list[str], takeaway: str) -> None:
    clear_slide(slide)
    add_background(slide)
    add_title(slide, title, subtitle)
    add_panel(slide, 0.55, 1.22, 8.4, 5.35)
    slide.shapes.add_picture(str(fig_path), Inches(0.7), Inches(1.35), width=Inches(8.1))
    add_bullets(slide, bullets, 9.05, 1.38, 3.5, 4.7)
    add_takeaway(slide, takeaway)


def find_slide_index(prs: Presentation, title: str) -> int | None:
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip() == title:
                return idx
    return None


def remove_slide(prs: Presentation, idx: int) -> None:
    slide_id_list = prs.slides._sldIdLst
    slides = list(slide_id_list)
    slide = slides[idx]
    prs.part.drop_rel(slide.rId)
    slide_id_list.remove(slide)


def move_last_slide_to(prs: Presentation, target_idx: int) -> None:
    slide_id_list = prs.slides._sldIdLst
    slides = list(slide_id_list)
    last = slides[-1]
    slide_id_list.remove(last)
    slide_id_list.insert(target_idx, last)


def main() -> None:
    prs = Presentation(DECK)

    rebuild_slide(
        prs.slides[5],
        "Where Augmentation Helps, With Writing-Level Granularity",
        "Each cell is one exact positive-case writing family, not an average over multiple writing versions.",
        FIG_PERF,
        [
            "Validation winners are concentrated in a few exact families, especially both/contrastive, both/structured, and both/uncertainty under joint-like modes.",
            "Learning is much cleaner: both and paper_only families become positive mainly under joint or joint_reasoning, while data_only stays weak.",
            "Granularity shows that 'augmentation helps' is too coarse; the exact writing family matters.",
        ],
        "The writing family is part of the treatment: averaging across them hides real heterogeneity.",
    )

    rebuild_slide(
        prs.slides[6],
        "Communication Shift Is Real, But It Does Not Buy Performance By Itself",
        "Here the feature shift is shown at exact writing-family granularity relative to the matched no-input baseline.",
        FIG_CHAT,
        [
            "On learning, most both and paper_only families push communication sharply upward; validation shifts are much smaller.",
            "Data_only is unstable and can even reverse direction, especially on learning.",
            "This shift is not enough on its own: corr(chat shift, ΔR²) is -0.01 on validation and only 0.26 on learning.",
        ],
        "A meaningful belief shift is not the same as a predictive gain.",
    )

    new_title = "Do The LLM Shifts Match What Matters In The Data?"
    existing = find_slide_index(prs, new_title)
    if existing is not None:
        rebuild_slide(
            prs.slides[existing],
            new_title,
            "Comparison against the manuscript's E-net feature-importance and SHAP analysis; this is a feature-level, not interaction-level, audit.",
            FIG_ENET,
            [
                "Yes for communication: it is the highest-importance E-net feature and also the largest aligned LLM shift.",
                "But several features look miscalibrated: peer incentive cost and group size are overindexed, while game length and contribution framing look underweighted.",
                "This first pass compares feature-level main effects; checking interaction mismatch would require a second pass against manuscript Figure 5.",
            ],
            "The augmented LLM is not random, but some of its biggest shifts still look misaligned with the data-driven model.",
        )
    else:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        rebuild_slide(
            slide,
            new_title,
            "Comparison against the manuscript's E-net feature-importance and SHAP analysis; this is a feature-level, not interaction-level, audit.",
            FIG_ENET,
            [
                "Yes for communication: it is the highest-importance E-net feature and also the largest aligned LLM shift.",
                "But several features look miscalibrated: peer incentive cost and group size are overindexed, while game length and contribution framing look underweighted.",
                "This first pass compares feature-level main effects; checking interaction mismatch would require a second pass against manuscript Figure 5.",
            ],
            "The augmented LLM is not random, but some of its biggest shifts still look misaligned with the data-driven model.",
        )
        move_last_slide_to(prs, 7)

    prs.save(DECK)
    print(DECK)


if __name__ == "__main__":
    main()
