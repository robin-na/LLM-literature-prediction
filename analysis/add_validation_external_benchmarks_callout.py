from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DECK = ROOT / "slides" / "prediction_results_deck.pptx"


def remove_existing(slide, marker: str) -> None:
    for shape in list(slide.shapes):
        if hasattr(shape, "text") and marker in shape.text:
            sp = shape.element
            sp.getparent().remove(sp)


def main() -> None:
    prs = Presentation(DECK)
    slide = prs.slides[1]

    marker = "External benchmarks:"
    remove_existing(slide, marker)

    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(8.95),
        Inches(6.10),
        Inches(3.55),
        Inches(0.54),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(245, 248, 252)
    box.line.color.rgb = RGBColor(190, 204, 219)

    tb = slide.shapes.add_textbox(
        Inches(9.08),
        Inches(6.16),
        Inches(3.30),
        Inches(0.42),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "External benchmarks: E-net RMSE 4.52, corr 0.773, DA 0.75, R² 0.599; noise ceiling R² 0.641."
    r.font.name = "Aptos"
    r.font.size = Pt(8.5)
    r.font.color.rgb = RGBColor(31, 41, 51)

    prs.save(DECK)
    print(DECK)


if __name__ == "__main__":
    main()
