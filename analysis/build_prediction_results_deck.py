from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PLOTS = ROOT / "plots"
SLIDES = ROOT / "slides"

BG = RGBColor(247, 245, 242)
TEXT = RGBColor(31, 41, 51)
MUTED = RGBColor(90, 102, 114)
NAVY = RGBColor(22, 50, 79)
TEAL = RGBColor(42, 157, 143)
ORANGE = RGBColor(244, 162, 97)
RED = RGBColor(231, 111, 81)
PANEL = RGBColor(255, 255, 255)


def add_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(11.9), Inches(0.65))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = NAVY

    slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.55),
        Inches(0.98),
        Inches(12.05),
        Inches(0.06),
    ).fill.solid()
    shape = slide.shapes[-1]
    shape.fill.fore_color.rgb = TEAL
    shape.line.fill.background()

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.05), Inches(11.8), Inches(0.3))
        p = sub_box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = "Aptos"
        r.font.size = Pt(10.5)
        r.font.color.rgb = MUTED


def add_panel(slide, x: float, y: float, w: float, h: float) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = RGBColor(222, 226, 230)


def add_picture(slide, image: Path, x: float, y: float, w: float | None = None, h: float | None = None) -> None:
    kwargs = {}
    if w is not None:
        kwargs["width"] = Inches(w)
    if h is not None:
        kwargs["height"] = Inches(h)
    slide.shapes.add_picture(str(image), Inches(x), Inches(y), **kwargs)


def add_bullets(slide, lines: list[str], x: float, y: float, w: float, h: float, accent: RGBColor = ORANGE) -> None:
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(253, 250, 245)
    panel.line.color.rgb = RGBColor(235, 221, 208)

    slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(0.12),
        Inches(h),
    ).fill.solid()
    accent_shape = slide.shapes[-1]
    accent_shape.fill.fore_color.rgb = accent
    accent_shape.line.fill.background()

    box = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.16), Inches(w - 0.3), Inches(h - 0.2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Aptos"
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT
        p.level = 0
        p.bullet = True
        p.space_after = Pt(5)


def add_takeaway(slide, text: str) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.65),
        Inches(6.72),
        Inches(12.0),
        Inches(0.42),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.85), Inches(6.79), Inches(11.7), Inches(0.24))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)

    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        prs.slide_width,
        Inches(0.22),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = TEAL
    band.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.75), Inches(1.05), Inches(11.2), Inches(1.6))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "LLM Prediction Of Treatment Effects:\nElicitation, Augmentation, And Model Dependence"
    r.font.name = "Aptos Display"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = NAVY

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(2.85), Inches(10.8), Inches(0.9))
    p = sub.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = (
        "Validation and learning-wave results, benchmarked against elastic net and a "
        "control-equals-treatment null."
    )
    r.font.name = "Aptos"
    r.font.size = Pt(16)
    r.font.color.rgb = TEXT

    add_bullets(
        slide,
        [
            "Baseline GPT-4.1 is stronger than many augmentation variants on validation.",
            "Reasoning helps more reliably than most augmentation tweaks.",
            "Augmentation changes implied CONFIG-level beliefs, not just error metrics.",
            "GPT-4.1-nano does not preserve the main pattern cleanly.",
        ],
        0.8,
        4.05,
        6.0,
        1.85,
        accent=TEAL,
    )

    add_picture(slide, PLOTS / "validation_r2_benchmark.png", 7.0, 3.7, w=5.5)


def add_single_figure_slide(
    prs: Presentation,
    title: str,
    subtitle: str,
    image: Path,
    bullets: list[str],
    takeaway: str,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, title, subtitle)
    add_panel(slide, 0.55, 1.2, 8.3, 5.35)
    add_picture(slide, image, 0.68, 1.32, w=8.0)
    add_bullets(slide, bullets, 9.05, 1.35, 3.55, 4.9)
    add_takeaway(slide, takeaway)


def add_two_figure_slide(
    prs: Presentation,
    title: str,
    subtitle: str,
    left_image: Path,
    right_image: Path,
    bullets: list[str],
    takeaway: str,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, title, subtitle)
    add_panel(slide, 0.55, 1.22, 6.0, 4.85)
    add_panel(slide, 6.78, 1.22, 6.0, 4.85)
    add_picture(slide, left_image, 0.7, 1.35, w=5.7)
    add_picture(slide, right_image, 6.93, 1.35, w=5.7)
    add_bullets(slide, bullets, 0.9, 6.12, 11.8, 0.95, accent=RED)
    add_takeaway(slide, takeaway)


def add_takeaway_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Takeaways", "What the slide deck should leave the audience with")

    add_bullets(
        slide,
        [
            "On validation, the canonical GPT-4.1 baseline remains hard to beat even after a large augmentation search over report, RAG, abstract, and prompt variants.",
            "Within unaugmented prompting, GPT-4.1 baseline_reasoning is the strongest simple setup; elastic net still outperforms the LLM baselines on R^2.",
            "Augmentation is only useful in selective regimes, mainly with both/paper_only inputs; data_only is consistently weak.",
            "Prompting and augmentation change the model's implied causal beliefs about CONFIGs, especially chat, rewardExists, and showOtherSummaries.",
            "GPT-4.1-nano should be treated as a separate regime: weaker validation baselines, 71 missing outputs, and materially different augmentation patterns.",
        ],
        0.9,
        1.55,
        11.5,
        4.9,
        accent=TEAL,
    )
    add_takeaway(slide, "Family-level patterns are more trustworthy than exact winner rankings.")


def build() -> Path:
    SLIDES.mkdir(exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)
    add_single_figure_slide(
        prs,
        "Benchmark Before Augmentation",
        "Validation R² uses the control-equals-treatment null as the denominator baseline.",
        PLOTS / "validation_r2_benchmark.png",
        [
            "GPT-4.1 baseline R² = 0.116; GPT-4.1 baseline_reasoning improves to 0.362.",
            "Elastic net reaches R² = 0.599; the noise ceiling is 0.641.",
            "All GPT-4.1-nano baselines are at or below zero on validation.",
        ],
        "Elastic net is much closer to the ceiling than any non-augmented LLM baseline.",
    )
    add_single_figure_slide(
        prs,
        "Main Validation Result",
        "This summary includes the larger report/RAG/abstract search space plus the newer positive-case prompt variants.",
        PLOTS / "validation_augmentation_search_summary.png",
        [
            "Report filters: median ΔRMSE +0.46, only 20% of variants beat baseline on point estimate.",
            "RAG filters: median ΔRMSE +0.24; Abstract filters: median +0.04.",
            "Positive-case prompt variants are worse on average: median ΔRMSE +0.89 and median ΔR² -0.25.",
        ],
        "Typical validation augmentation variants do not improve on the canonical GPT-4.1 baseline.",
    )
    add_single_figure_slide(
        prs,
        "Validation: Correlation And Directional Accuracy",
        "A few augmentation pockets improve correlation, but directional-accuracy gains are rare.",
        PLOTS / "validation_augmentation_search_summary_corr_da.png",
        [
            "Report, RAG, and abstract filter families all have negative median Δcorrelation on validation.",
            "Positive-case prompt variants are the one family with positive median Δcorrelation (+0.027), but their median Δdirectional accuracy is -0.10.",
            "Directional-accuracy gains are sparse across the whole validation search space.",
        ],
        "On validation, any upside is much easier to find in correlation than in directional accuracy.",
    )
    add_single_figure_slide(
        prs,
        "Where Augmentation Helps",
        "Grouped by input family and elicitation mode across validation and learning waves.",
        PLOTS / "crosswave_augmentation_performance_heatmap.png",
        [
            "Validation: gains are limited and mostly concentrated in correlation rather than RMSE/R².",
            "Learning: both + joint is the strongest GPT-4.1 region; paper_only can also help.",
            "Data_only is consistently weak across waves and metrics.",
        ],
        "Augmentation helps only in selective regimes, not as a general recipe.",
    )
    add_single_figure_slide(
        prs,
        "Augmentation Changes The Model's Priors",
        "Cell values show how augmentation changes the predicted treatment-effect contrast associated with each binary CONFIG.",
        PLOTS / "crosswave_config_augmentation_binary_delta_heatmap.png",
        [
            "Changes are substantive, not cosmetic: augmentation shifts the model's response to chat, rewardExists, and showOtherSummaries.",
            "These shifts depend on both input family and elicitation mode.",
            "Performance differences and belief changes move together; augmentation is changing what the model keys on.",
        ],
        "The important story is not only prediction accuracy, but also which causal patterns the model infers from the literature.",
    )
    add_two_figure_slide(
        prs,
        "Nano Is Not A Drop-In Substitute",
        "The same analysis pipeline on GPT-4.1-nano produces a materially different pattern.",
        PLOTS / "crosswave_augmentation_performance_heatmap.png",
        PLOTS / "crosswave_41nano_augmentation_performance_heatmap.png",
        [
            "GPT-4.1 cross-wave search has 0 missing outputs; GPT-4.1-nano has 71 missing outputs, including 50 missing learning cases for both_refined_joint_reasoning.",
            "Some coarse trends survive, especially data_only being weak.",
            "But the winning regions change and several learning-wave gains disappear under nano.",
        ],
        "Model choice changes the story; nano should be treated as a separate regime, not a cheaper replica.",
    )
    add_takeaway_slide(prs)
    add_single_figure_slide(
        prs,
        "Appendix: Cross-Wave Robustness",
        "Raw performance is not directly comparable across waves, so robustness is shown on normalized/ranked summaries.",
        PLOTS / "crosswave_robustness_scatter.png",
        [
            "GPT-4.1 cross-wave overall rank Spearman is 0.61: moderate robustness, not a full replication.",
            "The best families are more stable than the exact top-ranked variant.",
            "This is why the main deck emphasizes family-level conclusions over exact winners.",
        ],
        "Validation and learning are related but not interchangeable evaluation regimes.",
    )
    add_two_figure_slide(
        prs,
        "Appendix: Raw CONFIG Tendency Maps",
        "These show the model's implied treatment-effect tendencies before differencing against the matched no-input baseline.",
        PLOTS / "crosswave_config_tendency_binary_heatmap.png",
        PLOTS / "crosswave_config_tendency_continuous_heatmap.png",
        [
            "Binary features summarize sign and magnitude shifts for yes/no CONFIGs.",
            "Continuous features summarize correlations between predicted treatment effect and CONFIG intensity.",
            "The augmentation-delta slide in the main deck is the cleaner causal comparison; these are the raw tendency maps behind it.",
        ],
        "Raw tendency maps are useful for appendix detail, but the matched augmentation deltas carry the cleaner interpretation.",
    )

    out = SLIDES / "prediction_results_deck.pptx"
    prs.save(out)
    return out


if __name__ == "__main__":
    path = build()
    print(path)
