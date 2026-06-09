from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DECK = ROOT / "slides" / "prediction_results_deck.pptx"
from plot_paths import LEARNING_PLOTS

FIG = LEARNING_PLOTS / "learning_r2_benchmark.png"

BG = RGBColor(247, 245, 242)
TEXT = RGBColor(31, 41, 51)
MUTED = RGBColor(90, 102, 114)
NAVY = RGBColor(22, 50, 79)
TEAL = RGBColor(42, 157, 143)
PANEL = RGBColor(255, 255, 255)


def add_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(11.9), Inches(0.65))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = NAVY

    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.55),
        Inches(0.98),
        Inches(12.05),
        Inches(0.06),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.05), Inches(11.8), Inches(0.3))
        p = sub_box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = "Aptos"
        r.font.size = Pt(10.5)
        r.font.color.rgb = MUTED


def add_panel(slide, x: float, y: float, w: float, h: float) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = RGBColor(222, 226, 230)


def add_bullets(slide, lines: list[str], x: float, y: float, w: float, h: float) -> None:
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(253, 250, 245)
    panel.line.color.rgb = RGBColor(235, 221, 208)

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(0.12),
        Inches(h),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()

    box = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.15), Inches(w - 0.32), Inches(h - 0.2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Aptos"
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT
        p.bullet = True
        p.space_after = Pt(5)


def add_takeaway(slide, text: str) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.65),
        Inches(6.72),
        Inches(12.0),
        Inches(0.42),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.85), Inches(6.79), Inches(11.7), Inches(0.24))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER


def find_existing_slide_index(prs: Presentation, title: str) -> int | None:
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                txt = shape.text.strip()
                if txt:
                    texts.append(txt)
        if texts and texts[0] == title:
            return i
    return None


def remove_slide(prs: Presentation, idx: int) -> None:
    slide_id_list = prs.slides._sldIdLst
    slides = list(slide_id_list)
    slide = slides[idx]
    rel_id = slide.rId
    prs.part.drop_rel(rel_id)
    slide_id_list.remove(slide)


def move_last_slide_to(prs: Presentation, target_idx: int) -> None:
    slide_id_list = prs.slides._sldIdLst
    slides = list(slide_id_list)
    last = slides[-1]
    slide_id_list.remove(last)
    slide_id_list.insert(target_idx, last)


def main() -> None:
    prs = Presentation(DECK)
    title = "Benchmark Before Augmentation (Learning Wave)"
    existing = find_existing_slide_index(prs, title)
    if existing is not None:
        remove_slide(prs, existing)
        prs.save(DECK)
        prs = Presentation(DECK)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(
        slide,
        title,
        "Baseline-only comparison on the learning wave; no elastic net or noise ceiling on this slide.",
    )
    add_panel(slide, 0.55, 1.22, 8.4, 5.35)
    slide.shapes.add_picture(str(FIG), Inches(0.7), Inches(1.35), width=Inches(8.1))
    add_bullets(
        slide,
        [
            "GPT-4.1 is fairly flat across elicitation modes on learning; baseline_reasoning is best on RMSE, baseline_reasoning and baseline are best on directional accuracy/correlation tradeoffs.",
            "GPT-4.1-nano behaves differently: joint_reasoning is the strongest nano baseline on RMSE and directional accuracy.",
            "Unlike validation, the learning-wave ordering is not dominated by a single clear baseline variant.",
        ],
        9.05,
        1.38,
        3.5,
        4.7,
    )
    add_takeaway(
        slide,
        "On the learning wave, the baseline ordering changes and GPT-4.1-nano no longer behaves like a weaker copy of GPT-4.1.",
    )

    move_last_slide_to(prs, 2)
    prs.save(DECK)
    print(DECK)


if __name__ == "__main__":
    main()
